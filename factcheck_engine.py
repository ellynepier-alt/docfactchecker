import os, re, json, tempfile
from docx import Document
from docx.shared import RGBColor
from docx.oxml.ns import qn

SUPPORTED_EXTS = {'.txt', '.md', '.docx', '.pdf', '.pptx', '.png', '.jpg', '.jpeg'}

# Distinctive terms that indicate a document is actually about the DoC guideline subject
# matter. Used to gate the coverage map so unrelated documents (e.g., about something
# else entirely) don't get recommendations falsely checked off.
# Terms distinctive enough to this exact guideline that a single mention is sufficient
# evidence the document is on-topic.
STRONG_ANCHOR_TERMS = [
    'disorders of consciousness', 'vegetative state', 'unresponsive wakefulness',
    'minimally conscious', ' mcs+', ' mcs-', 'crs-r', 'amantadine', 'locked-in syndrome',
    'emcs', 'ptcs', 'confusional state', 'coma recovery scale', 'disability rating scale',
    'multidisciplinary rehabilitation', 'multidisciplinary team', 'neurorehabilitation', 'prognostication',
    'arousal facilitation', 'sternocleidomastoid', 'heterotopic ossification', 'deep pressure stimulation',
    'command following', 'sensory stimulation protocol', 'cognitive rehabilitation',
    'splinting and orthotic', 'neuromuscular retraining',
]

# Terms that are suggestive but common enough elsewhere that we require at least 2
# total hits (from either tier) before treating the document as on-topic.
WEAK_ANCHOR_TERMS = [
    'traumatic brain injury', ' tbi ', ' coma', ' mcs ', 'brain injury', 'consciousness',
    'intensivist', 'physiatrist', 'neurologist', 'rehabilitation team', 'behavioral evaluation',
    'arousal', 'sensory stimulation', 'contracture', 'psychosocial support', 'advance care planning',
    'discharge planning', 'rehabilitation intervention', 'sustained eye', 'behavioral responsiveness',
]

DOMAIN_ANCHOR_TERMS = STRONG_ANCHOR_TERMS + WEAK_ANCHOR_TERMS

# Precise phrase signatures per recommendation, used instead of loose single-word matching
# (which caused false "covered" hits on unrelated documents). A recommendation is only
# counted as covered if at least one genuinely distinctive phrase/pattern below is found.
REC_COVERAGE_SIGNATURES = {
    '1': [r'multidisciplinary rehabilitation', r'specialized rehabilitation', r'rehabilitation team', r'refer[^.]{0,40}rehabilitation', r'multidisciplinary team'],
    '2a': [r'standardized neurobehavioral assessment', r'standardized behavioral assessment', r'valid and reliable[^.]{0,30}assessment'],
    '2b': [r'serial standardized (?:neurobehavioral|behavioral)? ?assessment', r'serial assessment'],
    '2c': [r'increase arousal before', r'arousal before evaluation', r'diminished arousal'],
    '2d': [r'confound[^.]{0,30}diagnosis', r'treat conditions[^.]{0,30}confound', r'confounding condition'],
    '2e': [r'multimodal evaluation', r'multimodal assessment'],
    '2f': [r'frequent reevaluation', r'frequent neurobehavioral reevaluation'],
    '3': [r'universally poor prognosis', r'first 28 days', r'poor prognosis[^.]{0,30}28 days'],
    '4': [r'serial standardized behavioral evaluation', r'trajectory of recovery', r'establishing prognosis'],
    '5': [r'traumatic vs/uws', r'disability rating scale', r'\bdrs\b[^.]{0,30}(?:2-3 month|prognos)', r'spect scan', r'\bp300\b', r'eeg reactivity'],
    '6': [r'nontraumatic post-anoxic', r'post-anoxic vs/uws', r'crs-r', r'somatosensory evoked potential', r'\bseps?\b'],
    '7': [r'permanent vs\b', r'permanent vegetative state', r'chronic vs/uws', r'discontinue[^.]{0,30}terminology'],
    '8': [r'prognostic factors', r'individual outcomes vary', r'traumatic etiology', r'nontraumatic etiology'],
    '9': [r'goals of care', r'long-term disability', r'medical decision-making forms?', r'advance care planning'],
    '10': [r'chronic phase', r'prognostic counseling', r'prognosis discussion'],
    '11': [r'patient and family preferences', r'family preferences'],
    '12': [r'medical complications', r'systematic assessment[^.]{0,30}(?:prevention|complication)'],
    '13': [r'pain or suffering', r'pain and suffering', r'assess[^.]{0,30}\bpain\b'],
    '14': [r'\bamantadine\b'],
    '15': [r'nonvalidated treatment', r'non-validated treatment', r'limitations of evidence', r'evidence-based information'],
    '16': [r'pediatric[^.]{0,30}diagnostic', r'children with prolonged doc'],
    '17': [r'pediatric prognosis', r'children[^.]{0,30}prognosis', r'natural history[^.]{0,30}children'],
    '18': [r'pediatric therap', r'no established therapies for children'],
}


def load_kb(path):
    with open(path, 'r', encoding='utf-8') as f:
        return json.load(f)


_OCR_STATUS = {}


def check_ocr_engine_available():
    """Check once whether the Tesseract OCR engine is actually reachable (not just the Python
    wrapper). Cached after the first call. Used to distinguish 'OCR is broken' from 'this image
    genuinely has no text in it' — the two look identical (empty string) without this check."""
    if 'available' in _OCR_STATUS:
        return _OCR_STATUS['available']
    try:
        import pytesseract
        pytesseract.get_tesseract_version()
        _OCR_STATUS['available'] = True
    except Exception:
        _OCR_STATUS['available'] = False
    return _OCR_STATUS['available']


def ocr_image_bytes(image_bytes):
    """Run OCR on raw image bytes; return empty string on any failure (corrupt/unsupported image, no OCR engine, etc.)."""
    try:
        import pytesseract
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(image_bytes))
        return pytesseract.image_to_string(img).strip()
    except Exception:
        return ''


def get_docx_images(doc):
    """Return [{'alt': str|None, 'ocr_text': str}] for every image in the document,
    including images inside floating text boxes/diagrams (which live under
    wp:anchor rather than wp:inline)."""
    images = []
    for blip in doc.element.body.iter(qn('a:blip')):
        rId = blip.get(qn('r:embed'))
        if not rId:
            continue
        # Walk up to the wp:inline or wp:anchor container, which holds the docPr (alt text) sibling.
        container = blip
        alt = None
        for _ in range(10):
            container = container.getparent()
            if container is None:
                break
            tag = container.tag.split('}')[-1]
            if tag in ('inline', 'anchor'):
                alt = _xml_descendant_attr(container, 'docPr', 'descr') or _xml_descendant_attr(container, 'docPr', 'title')
                break
        ocr_text = ''
        try:
            part = doc.part.related_parts[rId]
            ocr_text = ocr_image_bytes(part.blob)
        except Exception:
            pass
        images.append({'alt': alt, 'ocr_text': ocr_text})
    return images


def get_pptx_images(prs):
    """Return [{'alt': str|None, 'ocr_text': str}] for every picture shape, recursing into groups."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    images = []

    def walk(shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                alt = _xml_descendant_attr(shape._element, 'cNvPr', 'descr')
                ocr_text = ''
                try:
                    ocr_text = ocr_image_bytes(shape.image.blob)
                except Exception:
                    pass
                images.append({'alt': alt, 'ocr_text': ocr_text})
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shape.shapes)

    for slide in prs.slides:
        walk(slide.shapes)
    return images


def get_pdf_images(path):
    """Return [str] of OCR'd text for every embedded image in a PDF (no alt-text concept in plain PDF)."""
    import fitz
    texts = []
    doc = fitz.open(path)
    for page in doc:
        for img in page.get_images(full=True):
            xref = img[0]
            try:
                base = doc.extract_image(xref)
                ocr_text = ocr_image_bytes(base['image'])
                if ocr_text:
                    texts.append(ocr_text)
            except Exception:
                continue
    return texts


def extract_docx_textboxes(doc):
    """Text inside floating text boxes (e.g., diagram/decision-tree shapes) lives in
    <w:txbxContent> elements nested inside drawings, which doc.paragraphs/doc.tables
    never reach. Walk the XML directly to pick these up.

    Word stores each such shape twice for compatibility: once as modern DrawingML
    (inside mc:Choice) and once as legacy VML (inside mc:Fallback), both with identical
    text. We skip the Fallback copy so content isn't double-counted."""
    chunks = []
    for node in doc.element.body.iter():
        tag = node.tag.split('}')[-1] if isinstance(node.tag, str) else ''
        if tag != 'txbxContent':
            continue
        ancestor = node.getparent()
        in_fallback = False
        for _ in range(15):
            if ancestor is None:
                break
            if ancestor.tag.split('}')[-1] == 'Fallback':
                in_fallback = True
                break
            ancestor = ancestor.getparent()
        if in_fallback:
            continue
        for p in node.findall('.//' + qn('w:p')):
            line = ''.join(t.text or '' for t in p.iter(qn('w:t')))
            if line.strip():
                chunks.append(line)
    return chunks


def extract_text(path):
    ext = os.path.splitext(path)[1].lower()
    if ext in ['.txt', '.md']:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    if ext == '.docx':
        doc = Document(path)
        chunks = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        chunks.append(cell.text)
        chunks.extend(extract_docx_textboxes(doc))
        chunks.extend(img['ocr_text'] for img in get_docx_images(doc) if img['ocr_text'])
        return '\n'.join(chunks)
    if ext == '.pdf':
        import fitz
        doc = fitz.open(path)
        chunks = [page.get_text() for page in doc]
        chunks.extend(get_pdf_images(path))
        return '\n'.join(chunks)
    if ext == '.pptx':
        from pptx import Presentation
        prs = Presentation(path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = ''.join(run.text for run in para.runs)
                        if text.strip():
                            chunks.append(text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                chunks.append(cell.text)
                if shape.has_chart:
                    continue
        chunks.extend(img['ocr_text'] for img in get_pptx_images(prs) if img['ocr_text'])
        return '\n'.join(chunks)
    if ext in ('.png', '.jpg', '.jpeg'):
        with open(path, 'rb') as f:
            return ocr_image_bytes(f.read())
    raise ValueError('Unsupported file type')


def context(text, start, end, width=160):
    lo = max(0, start - width); hi = min(len(text), end + width)
    out = re.sub(r'\s+', ' ', text[lo:hi]).strip()
    return ('...' if lo else '') + out + ('...' if hi < len(text) else '')


def add_flag(flags, kind, severity, matched, issue, rec, ctx):
    flags.append({'kind': kind, 'severity': severity, 'matched': matched, 'issue': issue, 'rec': rec or '', 'context': ctx})


def count_syllables(word):
    word = word.lower()
    vowels = 'aeiouy'
    count = 0
    prev_vowel = False
    for ch in word:
        is_vowel = ch in vowels
        if is_vowel and not prev_vowel:
            count += 1
        prev_vowel = is_vowel
    if word.endswith('e') and count > 1:
        count -= 1
    return max(1, count)


ACRONYM_GLOSSARY = {
    'VS/UWS': 'Vegetative State/Unresponsive Wakefulness Syndrome',
    'MCS': 'Minimally Conscious State',
    'DoC': 'Disorders of Consciousness',
    'CRS-R': 'Coma Recovery Scale-Revised',
    'DRS': 'Disability Rating Scale',
    'TBI': 'Traumatic Brain Injury',
    'SPECT': 'Single Photon Emission Computed Tomography',
    'PET': 'Positron Emission Tomography',
    'fMRI': 'functional Magnetic Resonance Imaging',
    'EEG': 'Electroencephalography',
    'ERP': 'Event-Related Potential',
    'SEP': 'Somatosensory Evoked Potential',
    'TMS': 'Transcranial Magnetic Stimulation',
    'PCI': 'Perturbational Complexity Index',
    'EMG': 'Electromyography',
    'MOLST': 'Medical Orders for Life-Sustaining Treatment',
}


def analyze_clarity(text):
    sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
    words = re.findall(r"[A-Za-z']+", text)
    num_sentences = max(1, len(sentences))
    num_words = max(1, len(words))
    syllables = sum(count_syllables(w) for w in words)

    avg_sentence_len = num_words / num_sentences
    avg_syllables_per_word = syllables / num_words

    flesch_reading_ease = 206.835 - 1.015 * avg_sentence_len - 84.6 * avg_syllables_per_word
    flesch_kincaid_grade = 0.39 * avg_sentence_len + 11.8 * avg_syllables_per_word - 15.59

    long_sentences = []
    for s in sentences:
        wc = len(re.findall(r"[A-Za-z']+", s))
        if wc > 30:
            snippet = s.strip()
            if len(snippet) > 220:
                snippet = snippet[:220] + '...'
            long_sentences.append({'text': snippet, 'word_count': wc})
    long_sentences = long_sentences[:5]

    undefined_acronyms = []
    for acr, expansion in ACRONYM_GLOSSARY.items():
        if re.search(r'\b' + re.escape(acr) + r'\b', text):
            key_words = [w for w in re.findall(r'[A-Za-z]+', expansion) if len(w) > 3]
            found_expansion = any(re.search(re.escape(w), text, re.IGNORECASE) for w in key_words) if key_words else False
            if not found_expansion:
                undefined_acronyms.append({'acronym': acr, 'expansion': expansion})

    passive_matches = re.findall(r'\b(?:is|are|was|were|be|been|being)\s+\w+ed\b', text, re.IGNORECASE)
    passive_count = len(passive_matches)
    passive_ratio = passive_count / num_sentences

    suggestions = []
    if flesch_kincaid_grade > 14:
        suggestions.append(
            f"This material reads at roughly a {flesch_kincaid_grade:.1f} grade level (college and above). "
            "Consider simplifying sentence structure and terminology, especially if patients or families will read it."
        )
    elif flesch_kincaid_grade > 10:
        suggestions.append(
            f"This material reads at roughly a {flesch_kincaid_grade:.1f} grade level (high school). "
            "Reasonable for a clinician audience, but may still be dense for family-facing materials."
        )

    if long_sentences:
        suggestions.append(
            f"{len(long_sentences)} sentence(s) exceed 30 words. Consider breaking these into shorter sentences for clarity."
        )

    if undefined_acronyms:
        names = ', '.join(a['acronym'] for a in undefined_acronyms[:6])
        suggestions.append(
            f"These clinical acronyms appear without being spelled out: {names}. "
            "Consider defining them on first use for readers unfamiliar with DoC terminology."
        )

    if passive_ratio > 0.5:
        suggestions.append(
            "This material relies heavily on passive voice (e.g., 'was performed' rather than 'the clinician performed'). "
            "Active voice is often clearer, especially for family-facing materials."
        )

    if not suggestions:
        suggestions.append('No major clarity issues detected. Language and structure appear reasonably accessible.')

    return {
        'flesch_reading_ease': round(flesch_reading_ease, 1),
        'flesch_kincaid_grade': round(flesch_kincaid_grade, 1),
        'avg_sentence_length': round(avg_sentence_len, 1),
        'long_sentences': long_sentences,
        'undefined_acronyms': undefined_acronyms,
        'passive_voice_count': passive_count,
        'suggestions': suggestions,
    }


# ---------------------------------------------------------------------------
# Proofreading: spelling, grammar (basic), plagiarism, missing references, conciseness
# ---------------------------------------------------------------------------
WORDY_PHRASES = {
    'in order to': 'to', 'due to the fact that': 'because', 'at this point in time': 'now',
    'a large number of': 'many', 'in the event that': 'if', 'with regard to': 'regarding',
    'in spite of the fact that': 'although', 'for the purpose of': 'to',
    'a majority of': 'most', 'in a timely manner': 'promptly', 'is able to': 'can',
    'has the ability to': 'can', 'in the process of': '(often can be cut)',
    'it is important to note that': '(often can be cut)', 'a number of': 'several',
    'on a daily basis': 'daily', 'in close proximity to': 'near',
}

_EXTRA_DOMAIN_WORDS = {
    'doc', 'docs', 'tbi', 'mcs', 'uws', 'vs', 'crs-r', 'crsr', 'drs', 'spect', 'pet', 'fmri',
    'eeg', 'erp', 'sep', 'seps', 'tms', 'pci', 'emg', 'ptcs', 'emcs', 'aan', 'acrm', 'nidilrr',
    'amantadine', 'physiatrist', 'intensivist', 'neurorehabilitation', 'prognostication',
    'neuroimaging', 'electrophysiologic', 'unresponsiveness', 'subacute',
}


def build_domain_whitelist(kb):
    words = set(_EXTRA_DOMAIN_WORDS)
    for rec in kb.get('recommendations', []):
        for token in re.findall(r"[A-Za-z']+", rec.get('topic', '') + ' ' + rec.get('text', '')):
            if len(token) > 2:
                words.add(token.lower())
    for term in kb.get('terminology_definitions', []):
        for token in re.findall(r"[A-Za-z']+", term.get('term', '') + ' ' + term.get('definition', '')):
            words.add(token.lower())
    for kf in kb.get('key_facts', []):
        for token in re.findall(r"[A-Za-z']+", kf.get('fact', '')):
            words.add(token.lower())
    return words


def check_spelling(text, whitelist):
    try:
        from spellchecker import SpellChecker
    except Exception:
        return []
    sc = SpellChecker()
    sc.word_frequency.load_words(whitelist)

    words = re.findall(r"[A-Za-z']+", text)
    counts = {}
    for w in words:
        if len(w) < 3 or w.isupper() or any(ch.isdigit() for ch in w):
            continue
        wl = w.lower()
        if wl in whitelist:
            continue
        counts[wl] = counts.get(wl, 0) + 1

    if not counts:
        return []
    unknown = sc.unknown(list(counts.keys()))
    results = []
    for w in sorted(unknown, key=lambda x: -counts[x])[:20]:
        results.append({'word': w, 'count': counts[w], 'suggestion': sc.correction(w)})
    return results


_GRAMMAR_ABBREVIATIONS = {
    'e.g.', 'i.e.', 'etc.', 'vs.', 'fig.', 'al.', 'approx.', 'vol.', 'pp.',
    'no.', 'jr.', 'sr.', 'st.', 'dr.', 'mr.', 'mrs.', 'ms.', 'prof.', 'inc.', 'ltd.',
}


def check_grammar(text):
    issues = []
    for line in text.split('\n'):
        for m in re.finditer(r'\b(\w+)\s+\1\b', line, re.IGNORECASE):
            issues.append({'issue': f'Repeated word: "{m.group(1)}"', 'context': line[max(0, m.start() - 40):m.end() + 40].strip()})
    if '  ' in text:
        first = text.find('  ')
        issues.append({'issue': 'Multiple consecutive spaces found.', 'context': text[max(0, first - 30):first + 30].strip()})
    for m in re.finditer(r'(\S+)([.!?])\s+([a-z])', text):
        preceding_token = m.group(1)
        combined = (preceding_token + m.group(2)).lower()
        if any(combined.endswith(ab) for ab in _GRAMMAR_ABBREVIATIONS):
            continue
        if re.fullmatch(r'\(?[0-9]{1,3}\)?|\(?[a-z]\)?|\(?[ivxlcdm]{1,4}\)?', preceding_token, re.IGNORECASE):
            continue  # likely a list marker like "1." "a." "iv." — not a real sentence boundary
        following_context = text[m.end() - 1:m.end() + 15]
        if re.match(r'\s*(?:https?://|www\.|doi\.org|[a-z0-9.]+\.(?:com|org|edu|gov)\b)', following_context, re.IGNORECASE):
            continue  # URL/DOI right after the period, common in reference lists
        issues.append({'issue': 'Sentence may need to start with a capital letter.', 'context': text[max(0, m.start() - 20):m.start() + 25].strip()})
    if text.count('(') != text.count(')'):
        issues.append({'issue': 'Unmatched parentheses found in the document.', 'context': ''})
    if text.count('"') % 2 != 0:
        issues.append({'issue': 'Unmatched double-quote mark found in the document.', 'context': ''})
    return issues[:20]


def check_plagiarism(text, kb):
    ref_sentences = []
    for rec in kb.get('recommendations', []):
        ref_sentences.append((f"Recommendation {rec['id']}", rec['text']))
    for term in kb.get('terminology_definitions', []):
        ref_sentences.append((term['term'], term['definition']))

    def ngrams(s, n=8):
        w = re.findall(r"[a-z']+", s.lower())
        return set(tuple(w[i:i + n]) for i in range(len(w) - n + 1)) if len(w) >= n else set()

    ref_ngrams = [(label, ngrams(rt)) for label, rt in ref_sentences]
    doc_sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]

    findings = []
    for ds in doc_sentences:
        dg = ngrams(ds)
        if not dg:
            continue
        for label, rg in ref_ngrams:
            if dg & rg:
                has_quotes = '"' in ds or '\u201c' in ds or '\u201d' in ds or "'" in ds
                findings.append({'matched_source': label, 'sentence': ds[:220], 'quoted': has_quotes})
                break
    return findings[:10]


def check_internal_duplication(text):
    paras = [p.strip() for p in re.split(r'\n', text) if len(p.strip()) > 40]
    seen = {}
    dups = []
    for p in paras:
        key = re.sub(r'\s+', ' ', p.lower())
        if key in seen:
            dups.append(p[:180])
        else:
            seen[key] = True
    return dups[:10]


_CLAIM_INDICATOR = re.compile(
    r'\b(?:studies show|research shows|shown to|proven to|according to|evidence suggests|data indicates|'
    r'studies indicate|clinical trials|meta-analysis|systematic review|significantly (?:higher|lower|increases|decreases|improves))\b',
    re.IGNORECASE)
_STAT_PATTERN = re.compile(r'\b\d{1,3}(?:\.\d+)?\s?%|\bp\s*[<=]\s*0\.\d+')
_CITATION_PATTERN = re.compile(r'\(\s*[A-Z][a-zA-Z]+(?:\s+et al\.?)?,?\s*\d{4}\s*\)|\[\d+\]|https?://|doi\.org')


def check_missing_references(text):
    sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
    flagged = []
    for s in sentences:
        if (_CLAIM_INDICATOR.search(s) or _STAT_PATTERN.search(s)) and not _CITATION_PATTERN.search(s):
            flagged.append(s[:220])
    return flagged[:10]


def check_conciseness(text):
    low = text.lower()
    findings = []
    for phrase, suggestion in WORDY_PHRASES.items():
        count = low.count(phrase)
        if count:
            findings.append({'phrase': phrase, 'count': count, 'suggestion': suggestion})
    return findings


def strip_references_section(text):
    """Reference lists/bibliographies have irregular formatting (URLs, journal abbreviations,
    page ranges) that trips grammar/spelling/citation checks meant for prose. Cut everything
    from a References/Bibliography/Works Cited heading onward."""
    m = re.search(r'^[ \t]*(references|bibliography|works cited|reference list|citations)[ \t]*:?[ \t]*$', text, re.IGNORECASE | re.MULTILINE)
    if m:
        return text[:m.start()]
    return text


def analyze_proofreading(text, kb):
    whitelist = build_domain_whitelist(kb)
    body_text = strip_references_section(text)
    return {
        'spelling': check_spelling(body_text, whitelist),
        'grammar': check_grammar(body_text),
        'plagiarism': check_plagiarism(body_text, kb),
        'internal_duplication': check_internal_duplication(body_text),
        'missing_references': check_missing_references(body_text),
        'conciseness': check_conciseness(body_text),
    }


def _xml_descendant_attr(element, local_tag, attr):
    """Find first descendant with the given local (namespace-stripped) tag and return an attribute."""
    for node in element.iter():
        tag = node.tag.split('}')[-1] if isinstance(node.tag, str) else ''
        if tag == local_tag:
            return node.get(attr)
    return None


def assess_image_alt_quality(images):
    """images: list of {'alt': str|None, 'ocr_text': str}. Returns (total, missing, poor, feedback[])."""
    generic_terms = {'image', 'picture', 'photo', 'graphic', 'img', 'untitled', 'diagram', 'chart', 'picture1', 'graphic1'}
    total = len(images)
    missing = 0
    poor = 0
    feedback = []
    for im in images:
        alt = (im.get('alt') or '').strip()
        ocr = (im.get('ocr_text') or '').strip()
        if not alt:
            missing += 1
            if ocr:
                snippet = ocr[:150].replace('\n', ' ')
                feedback.append(f'Missing alt text on an image that contains readable text: "{snippet}" — consider using this (or a summary of it) as the alt text.')
            else:
                feedback.append('Missing alt text on an image with no machine-readable text detected — add a description of what the image depicts.')
        elif alt.lower().strip('.') in generic_terms or len(alt) < 4:
            poor += 1
            feedback.append(f'Alt text "{alt}" is too generic to convey meaning to screen-reader users.')
        elif ocr and len(ocr) > 15:
            ocr_words = set(re.findall(r'[a-z]{4,}', ocr.lower()))
            alt_words = set(re.findall(r'[a-z]{4,}', alt.lower()))
            if ocr_words and not (ocr_words & alt_words):
                poor += 1
                preview = ocr[:100].replace('\n', ' ')
                feedback.append(f'Alt text "{alt}" does not appear to reflect the visible text in the image ("{preview}..."). Screen-reader users may miss this information entirely.')
    return total, missing, poor, feedback


def check_docx_link_text(doc):
    generic_texts = {'click here', 'here', 'more', 'read more', 'link', 'this link', 'click', 'learn more', 'more info', 'info'}
    total_links = 0
    bad_links = []
    for hyperlink in doc.element.body.iter(qn('w:hyperlink')):
        total_links += 1
        text = ''.join(t.text or '' for t in hyperlink.iter(qn('w:t'))).strip()
        if text.lower() in generic_texts:
            bad_links.append(text)
    return total_links, bad_links


def check_accessibility_docx(path):
    findings = []
    doc = Document(path)

    images = get_docx_images(doc)
    total_images, missing_alt, poor_alt, alt_feedback = assess_image_alt_quality(images)
    if total_images:
        compliant_images = total_images - missing_alt - poor_alt
        pct = round(100 * compliant_images / total_images)
        findings.append({
            'check': 'Image alternative text (OCR-verified)', 'wcag': '1.1.1 Non-text Content (Level A)',
            'category': 'Images', 'percent': pct,
            'status': 'fail' if (missing_alt or poor_alt) else 'pass',
            'detail': f'{compliant_images} of {total_images} image(s) ({pct}%) have adequate alt text; {missing_alt} missing, {poor_alt} generic or mismatched with the image\'s actual visible content (checked via OCR).',
            'fix': ('Right-click the image → "Edit Alt Text..." (or Format Picture → Alt Text pane) → write a concise description of what the image '
                    'shows or the information it conveys. If the image is purely decorative and adds no information, mark it as decorative instead '
                    'of leaving the field blank.'),
        })
        for fb in alt_feedback[:8]:
            findings.append({
                'check': 'Alt text feedback', 'wcag': '1.1.1 Non-text Content (Level A)', 'status': 'warn', 'detail': fb,
                'fix': 'Right-click the image → "Edit Alt Text..." → replace the current text with an accurate description of the image\'s content or purpose.',
            })
    else:
        findings.append({'check': 'Image alternative text', 'wcag': '1.1.1 Non-text Content (Level A)', 'category': 'Images', 'percent': None, 'status': 'na', 'detail': 'No images found in this document.'})

    heading_used = any(p.style and p.style.name.startswith('Heading') and p.text.strip() for p in doc.paragraphs)
    real_headings = sum(1 for p in doc.paragraphs if p.style and p.style.name.startswith('Heading') and p.text.strip())
    fake_headings = 0
    for p in doc.paragraphs:
        if p.style and p.style.name.startswith('Heading'):
            continue
        txt = p.text.strip()
        if txt and len(txt) < 80 and p.runs and all(r.bold for r in p.runs if r.text.strip()):
            fake_headings += 1
    heading_total = real_headings + fake_headings
    heading_pct = round(100 * real_headings / heading_total) if heading_total else (100 if heading_used else None)
    findings.append({
        'check': 'Heading styles used for structure', 'wcag': '1.3.1 Info and Relationships (A) / 2.4.6 Headings and Labels (AA)',
        'category': 'Headings & structure', 'percent': heading_pct,
        'status': 'pass' if heading_used else 'fail',
        'detail': ('Document uses Word Heading styles, which screen readers rely on for section navigation.' if heading_used
                   else 'No paragraphs use Word Heading styles, so screen-reader users cannot navigate by section.'),
        'fix': ('Select each section title → go to the Home tab → choose "Heading 1", "Heading 2", etc. from the Styles gallery, instead of manually '
                'bolding or enlarging the font. This lets screen-reader users jump between sections and lets Word auto-generate a table of contents.'),
    })
    if fake_headings:
        findings.append({
            'check': 'Bold text used in place of headings', 'wcag': '1.3.1 Info and Relationships (Level A)',
            'status': 'warn', 'detail': f'{fake_headings} short bold line(s) look like section titles but are not tagged with a Heading style, so they are invisible to screen-reader navigation.',
            'fix': 'Select each bolded title line → apply a Heading style from the Home tab\'s Styles gallery (e.g., Heading 1 or Heading 2) instead of just bold text.',
        })

    tables_total = len(doc.tables)
    tables_missing_header = 0
    tables_visual_only = 0
    for t in doc.tables:
        has_header_flag = False
        if t.rows:
            tr = t.rows[0]._tr
            trPr = tr.find(qn('w:trPr'))
            if trPr is not None and trPr.find(qn('w:tblHeader')) is not None:
                has_header_flag = True
        if not has_header_flag:
            tables_missing_header += 1
            tblPr = t._tbl.find(qn('w:tblPr'))
            tblLook = tblPr.find(qn('w:tblLook')) if tblPr is not None else None
            first_row_style = tblLook is not None and tblLook.get(qn('w:firstRow')) == '1'
            if first_row_style:
                tables_visual_only += 1
    if tables_total:
        tables_compliant = tables_total - tables_missing_header
        pct = round(100 * tables_compliant / tables_total)
        detail = f'{tables_compliant} of {tables_total} table(s) ({pct}%) have a designated header row; {tables_missing_header} do not, so screen readers cannot announce column context for those data cells.'
        fix = 'Click into the header row → Table Properties (Layout tab → Properties, or right-click → Table Properties) → Row tab → check "Repeat as header row at the top of each page."'
        if tables_visual_only:
            detail += (f' Note: {tables_visual_only} of the non-compliant table(s) use the "Header Row" table-style option, which only changes visual '
                       f'formatting (bold/shading) — it does NOT mark the row as an accessible header.')
            fix = ('The bold first row you see comes from the "Header Row" table-STYLE option (Table Design tab), which only affects appearance. '
                   'To actually tag it as a header: click into that row → Table Properties → Row tab → check "Repeat as header row at the top of each page." '
                   'Do this for each table.')
        findings.append({
            'check': 'Table header rows', 'wcag': '1.3.1 Info and Relationships (Level A)',
            'category': 'Tables', 'percent': pct,
            'status': 'fail' if tables_missing_header else 'pass',
            'detail': detail,
            'fix': fix,
        })
    else:
        findings.append({'check': 'Table header rows', 'wcag': '1.3.1 Info and Relationships (Level A)', 'category': 'Tables', 'percent': None, 'status': 'na', 'detail': 'No tables found in this document.'})

    total_links, bad_links = check_docx_link_text(doc)
    if total_links:
        good_links = total_links - len(bad_links)
        pct = round(100 * good_links / total_links)
        findings.append({
            'check': 'Meaningful link text', 'wcag': '2.4.4 Link Purpose in Context (Level A)',
            'category': 'Links', 'percent': pct,
            'status': 'fail' if bad_links else 'pass',
            'detail': (f'{good_links} of {total_links} link(s) ({pct}%) use descriptive text; {len(bad_links)} use generic text like "{bad_links[0]}" that gives no context out of place — screen-reader users often navigate by a list of links alone.' if bad_links
                       else f'All {total_links} link(s) ({pct}%) use descriptive text.'),
            'fix': 'Right-click the link → "Edit Hyperlink..." → change the displayed text to describe the destination (e.g., "CDC brain injury guidelines" instead of "click here").',
        })
    else:
        findings.append({'check': 'Meaningful link text', 'wcag': '2.4.4 Link Purpose in Context (Level A)', 'category': 'Links', 'percent': None, 'status': 'na', 'detail': 'No hyperlinks found in this document.'})

    return findings


def check_accessibility_pptx(path):
    from pptx import Presentation

    findings = []
    prs = Presentation(path)
    slides = list(prs.slides)

    slides_without_title = 0
    for slide in slides:
        has_title = slide.shapes.title is not None and slide.shapes.title.has_text_frame and slide.shapes.title.text_frame.text.strip()
        if not has_title:
            slides_without_title += 1

    images = get_pptx_images(prs)
    total_images, missing_alt, poor_alt, alt_feedback = assess_image_alt_quality(images)
    if total_images:
        compliant_images = total_images - missing_alt - poor_alt
        pct = round(100 * compliant_images / total_images)
        findings.append({
            'check': 'Image alternative text (OCR-verified)', 'wcag': '1.1.1 Non-text Content (Level A)',
            'category': 'Images', 'percent': pct,
            'status': 'fail' if (missing_alt or poor_alt) else 'pass',
            'detail': f'{compliant_images} of {total_images} image(s) ({pct}%) have adequate alt text; {missing_alt} missing, {poor_alt} generic or mismatched with the image\'s actual visible content (checked via OCR).',
            'fix': 'Right-click the image → "Edit Alt Text..." → write a concise description of what the image shows or the information it conveys. Mark purely decorative images as decorative instead.',
        })
        for fb in alt_feedback[:8]:
            findings.append({
                'check': 'Alt text feedback', 'wcag': '1.1.1 Non-text Content (Level A)', 'status': 'warn', 'detail': fb,
                'fix': 'Right-click the image → "Edit Alt Text..." → replace the current text with an accurate description of the image\'s content or purpose.',
            })
    else:
        findings.append({'check': 'Image alternative text', 'wcag': '1.1.1 Non-text Content (Level A)', 'category': 'Images', 'percent': None, 'status': 'na', 'detail': 'No images found in this presentation.'})

    slides_with_title = len(slides) - slides_without_title
    slide_pct = round(100 * slides_with_title / len(slides)) if slides else None
    findings.append({
        'check': 'Slide titles', 'wcag': '2.4.6 Headings and Labels (AA) / 1.3.1 Info and Relationships (A)',
        'category': 'Slide titles', 'percent': slide_pct,
        'status': 'fail' if slides_without_title else 'pass',
        'detail': f'{slides_with_title} of {len(slides)} slide(s) ({slide_pct}%) have a title placeholder; {slides_without_title} do not, so screen readers cannot announce the topic of those slides.',
        'fix': ('Use the slide layout\'s built-in Title placeholder ("Click to add title") rather than a plain text box. If a slide is missing a title '
                'placeholder, go to Home → Layout and choose a layout that includes one, then fill it in — you can format/hide it visually if needed, '
                'but keep it present for screen readers.'),
    })

    return findings


def check_accessibility_pdf(path):
    findings = [
        {
            'check': 'Tagged PDF structure', 'wcag': '1.3.1 Info and Relationships (A) / 4.1.2 Name, Role, Value (A)',
            'status': 'manual',
            'detail': "Automatic tag detection isn't available in this tool. Verify this PDF is tagged (e.g., with Acrobat's Accessibility Checker) so screen readers can interpret headings, tables, and reading order.",
            'fix': ('In Adobe Acrobat Pro: Tools → Accessibility → "Autotag Document" to generate initial tags, then run the Accessibility Checker '
                    '(Tools → Accessibility → Full Check) and manually fix any headings, tables, or reading-order issues it flags. If the PDF was '
                    'exported from Word/PowerPoint, it\'s usually far easier to fix the accessibility issues in the original file and re-export/re-tag from there.'),
        },
    ]
    image_texts = get_pdf_images(path)
    if image_texts:
        preview = image_texts[0][:150].replace('\n', ' ')
        findings.append({
            'check': 'Image alternative text', 'wcag': '1.1.1 Non-text Content (Level A)',
            'status': 'manual',
            'detail': f'{len(image_texts)} image(s) contain machine-readable text (e.g., "{preview}..."). This text has been included in the fact-check, but PDF alt-text tagging still needs manual verification (e.g., with Acrobat\'s Accessibility Checker).',
            'fix': 'In Acrobat Pro: Tools → Accessibility → Reading Order → select each image → "Add Alternate Text," or use the Accessibility Checker\'s "Add Alternate Text" flow.',
        })
    else:
        findings.append({
            'check': 'Image alternative text', 'wcag': '1.1.1 Non-text Content (Level A)',
            'status': 'manual',
            'detail': "No machine-readable text was detected in this PDF's images. Alt-text tagging still can't be reliably verified automatically — check with Acrobat's Accessibility Checker.",
            'fix': 'In Acrobat Pro: Tools → Accessibility → Reading Order → select each image → "Add Alternate Text," or use the Accessibility Checker\'s "Add Alternate Text" flow.',
        })
    return findings


def check_accessibility(path, text):
    ext = os.path.splitext(path)[1].lower()
    if ext == '.docx':
        findings = check_accessibility_docx(path)
    elif ext == '.pptx':
        findings = check_accessibility_pptx(path)
    elif ext == '.pdf':
        findings = check_accessibility_pdf(path)
    elif ext in ('.png', '.jpg', '.jpeg'):
        findings = [{
            'check': 'Standalone image alternative text', 'wcag': '1.1.1 Non-text Content (Level A)', 'status': 'manual',
            'detail': ('This is a standalone image file, so alt text isn\'t part of the file itself — it depends on where the image is used '
                       '(e.g., the alt attribute when embedded on a webpage, or alt text in a document). Text visible in the image has been '
                       'extracted via OCR and included in the fact-check above.'),
        }]
    else:
        findings = [{
            'check': 'Format limitations', 'wcag': 'N/A', 'status': 'na',
            'detail': 'Plain text/Markdown files have no document-structure accessibility concerns beyond readability (see Clarity section above).',
        }]
    return findings


# Only these "primary" checks count toward the compliance score. Supplementary
# per-item detail (e.g. individual "Alt text feedback" entries, "Bold text used
# in place of headings") is excluded so one bad image doesn't get double-penalized.
SCORED_ACCESSIBILITY_CHECKS = {
    'Image alternative text', 'Image alternative text (OCR-verified)',
    'Heading styles used for structure', 'Table header rows',
    'Meaningful link text', 'Slide titles', 'Tagged PDF structure',
}


def compute_accessibility_score(findings):
    categories = []
    for f in findings:
        if f['check'] not in SCORED_ACCESSIBILITY_CHECKS:
            continue
        pct = f.get('percent')
        if pct is None:
            # Fallback for checks without a granular percent (e.g. tagged-PDF, which is
            # inherently binary/manual) — treat pass/fail as 100/0, skip na/manual entirely.
            if f['status'] == 'pass':
                pct = 100
            elif f['status'] == 'fail':
                pct = 0
            elif f['status'] == 'warn':
                pct = 50
            else:
                continue
        categories.append({'category': f.get('category', f['check']), 'percent': pct})

    if not categories:
        return {'score': None, 'zone': None, 'breakdown': []}

    score = round(sum(c['percent'] for c in categories) / len(categories))
    if score >= 90:
        zone = 'green'
    elif score >= 70:
        zone = 'yellow'
    else:
        zone = 'red'
    return {'score': score, 'zone': zone, 'breakdown': categories}


def detect_quiz_spans(norm):
    """Find quiz-style blocks (question stem + lettered/numbered options + an Answer: marker) so
    their content isn't treated as the document's own assertions. Returns list of (start, end, answer_letter).

    Each question's stem is bounded by the end of the PREVIOUS question's span (so consecutive
    questions can't bleed into each other) and capped at a max lookback distance (so a quiz late
    in a long document doesn't swallow unrelated preceding prose)."""
    spans = []
    prev_end = 0
    MAX_STEM_LOOKBACK = 600
    for m in re.finditer(r'\banswer\s*[:\-]?\s*([a-e]|\d{1,2})\b', norm, re.IGNORECASE):
        region_start = prev_end
        region = norm[region_start:m.start()]
        options = list(re.finditer(r'(?:(?<=[\s.])|^)(?:[a-e][\.\)]|\d{1,2}\))\s', region, re.IGNORECASE))
        if len(options) < 2:
            continue  # not enough lettered/numbered options nearby — probably not a real quiz block
        first_option_pos = region_start + options[0].start()
        stem_start = max(region_start, first_option_pos - MAX_STEM_LOOKBACK)
        block_end = m.end()
        nxt_period = norm.find('.', block_end)
        block_end = nxt_period + 1 if (nxt_period != -1 and nxt_period - block_end < 300) else min(len(norm), block_end + 100)
        spans.append((stem_start, block_end, m.group(1).upper()))
        prev_end = block_end
    return spans


def parse_quiz_options(block_text):
    """Map each option letter to its own text, e.g. {'A': '28 days', 'B': '3 months', ...}.
    Letters only (not numbered 1/2/3) — numbered matching corrupted parsing whenever an option's
    own text ended in a number (e.g., "a total score of 2", "(Score: 2)"), which is extremely common."""
    options = {}
    matches = list(re.finditer(r'\b([a-e])[\.\)]\s*', block_text, re.IGNORECASE))
    for i, m in enumerate(matches):
        start = m.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(block_text)
        options[m.group(1).upper()] = block_text[start:end].strip(' .')
    return options


def verify_quiz_answer(block_text, answer_letter, kb):
    """Best-effort check of a quiz's stated correct answer against known guideline numeric facts
    (3 vs 12 month rule, amantadine dose, CRS-R cutoff). Returns a list of finding dicts."""
    findings = []
    ans_marker = re.search(r'\banswer\b', block_text, re.IGNORECASE)
    options_region = block_text[:ans_marker.start()] if ans_marker else block_text
    options = parse_quiz_options(options_region)
    chosen = options.get(answer_letter, '')
    if not chosen:
        return findings

    first_option_pos = min((m.start() for m in re.finditer(r'\b[a-e][\.\)]', options_region, re.IGNORECASE)), default=len(options_region))
    stem = options_region[:first_option_pos]

    if ('permanent' in stem or 'discontinu' in stem) and ('vegetative' in stem or ' vs' in stem or 'uws' in stem):
        nums = re.findall(r'(\d{1,2})\s*months?', chosen)
        if nums:
            stated = nums[0]
            stem_norm = stem.replace('non-traumatic', 'nontraumatic')
            if 'nontraumatic' in stem_norm and 'traumatic' not in stem_norm.replace('nontraumatic', ''):
                if stated != '3':
                    findings.append({'issue': f"Quiz answer states {stated} months for nontraumatic VS/UWS, but the guideline specifies 3 months.", 'rec': '7'})
            elif 'traumatic' in stem_norm:
                if stated != '12':
                    findings.append({'issue': f"Quiz answer states {stated} months for traumatic VS/UWS, but the guideline specifies 12 months.", 'rec': '7'})

    if 'amantadine' in stem:
        m = re.search(r'(\d{2,4})\s*mg', chosen)
        if m and m.group(1) not in ('100', '200'):
            findings.append({'issue': f"Quiz answer states {m.group(1)} mg for amantadine, but the guideline specifies 100-200 mg twice daily.", 'rec': '14'})

    if 'crs-r' in stem:
        subscale_keywords = ('auditory', 'visual', 'motor', 'oromotor', 'verbal', 'communication', 'arousal')
        subscale_hits = sum(1 for kw in subscale_keywords if kw in chosen.lower())
        if subscale_hits < 2:  # not a subscale-breakdown answer, so a bare number plausibly means a total/overall score
            m = re.search(r'(?:total\s+)?score[^.]{0,20}?(?:of\s*)?(\d{1,2})\b', chosen, re.IGNORECASE)
            if m and int(m.group(1)) < 6:
                findings.append({'issue': f"Quiz answer implies a CRS-R score of {m.group(1)} is favorable, but the guideline associates scores of 6 or higher with increased likelihood of recovery.", 'rec': '6'})

    return findings


def _flag_position(norm, flag):
    matched = flag.get('matched', '')
    if matched:
        matched_norm = re.sub(r'\s+', ' ', matched.lower()).strip()
        if matched_norm:
            pos = norm.find(matched_norm)
            if pos != -1:
                return pos
    ctx = flag.get('context', '')
    if not ctx:
        return -1
    if ctx.startswith('...'):
        ctx = ctx[3:]
    if ctx.endswith('...'):
        ctx = ctx[:-3]
    ctx_norm = re.sub(r'\s+', ' ', ctx.lower()).strip()
    return norm.find(ctx_norm)


def compute_ocr_warning(filepath, text):
    ext = os.path.splitext(filepath)[1].lower()
    image_dependent_exts = ('.png', '.jpg', '.jpeg')
    image_containing_exts = ('.docx', '.pptx', '.pdf')

    if ext in image_dependent_exts:
        if not check_ocr_engine_available():
            return ("The OCR engine (Tesseract) isn't available in this environment, so no text could be extracted from this image. "
                    "Everything below will be empty or inaccurate — this image was NOT actually fact-checked. "
                    "Contact your administrator to verify Tesseract is installed (see packages.txt).")
        if not text.strip():
            return ("No readable text was detected in this image via OCR. If the image does contain text, try a higher-resolution "
                    "version — otherwise this image was not fact-checked because there was nothing to check.")
    elif ext in image_containing_exts:
        if not check_ocr_engine_available():
            return ("The OCR engine (Tesseract) isn't available in this environment, so text inside any images embedded in this "
                    "document (charts, diagrams, screenshots) was NOT extracted or fact-checked — only the document's regular "
                    "text was checked. Contact your administrator to verify Tesseract is installed (see packages.txt).")
    return None


def run_checks(filepath, kb):
    text = extract_text(filepath)
    ocr_warning = compute_ocr_warning(filepath, text)
    flags = []
    low = text.lower()
    norm = re.sub(r'[ \t]*\n+[ \t]*', '. ', low)
    norm = re.sub(r'\s+', ' ', norm)
    norm = re.sub(r'\.\s*\.', '.', norm)

    for f in kb['terminology_flags']:
        for m in re.finditer(re.escape(f['pattern'].lower()), low):
            add_flag(flags, 'Terminology', f['severity'], text[m.start():m.end()], f['issue'], f.get('rec'), context(text, m.start(), m.end()))

    for f in kb['contradiction_flags']:
        for m in re.finditer(re.escape(f['pattern'].lower()), low):
            add_flag(flags, 'Possible contradiction', f['severity'], text[m.start():m.end()], f['issue'], f.get('rec'), context(text, m.start(), m.end()))

    if 'amantadine' in norm:
        for m in re.finditer(r'amantadine[^.]{0,80}?(\d{2,4})\s*mg', norm):
            dose = m.group(1)
            if dose not in ['100', '200']:
                add_flag(flags, 'Key fact mismatch', 'high', f'amantadine {dose} mg', 'Guideline dose is amantadine 100-200 mg twice daily.', '14', context(norm, m.start(), m.end()))
        # (Amantadine mention no longer flagged for review — only explicit dose/window mismatches above are flagged.)

    for m in re.finditer(r'(?<!non)(?<!non-)traumatic[^.]{0,90}?3\s*month', norm):
        add_flag(flags, 'Key fact mismatch', 'high', 'traumatic ... 3 months', 'Three months applies to nontraumatic VS/UWS; traumatic VS/UWS uses 12 months.', '7', context(norm, m.start(), m.end()))
    for m in re.finditer(r'non-?traumatic[^.]{0,90}?12\s*month', norm):
        add_flag(flags, 'Key fact mismatch', 'high', 'nontraumatic ... 12 months', 'Twelve months applies to traumatic VS/UWS; nontraumatic VS/UWS uses 3 months.', '7', context(norm, m.start(), m.end()))

    if 'crs-r' in norm:
        for m in re.finditer(r'crs-r[^.]{0,60}?score[^.]{0,20}?(?:of\s*)?(\d{1,2})', norm):
            score = int(m.group(1))
            if score < 6:
                add_flag(flags, 'Key fact mismatch', 'high', f'CRS-R score {score}', 'Guideline associates CRS-R scores of 6 or higher (>1 month after onset) with increased likelihood of recovery in nontraumatic post-anoxic VS/UWS.', '6', context(norm, m.start(), m.end()))

    for m in re.finditer(r'\bdrs\b[^.]{0,80}?(\d{1,2})\s*(?:-|to)\s*(\d{1,2})\s*month', norm):
        lo, hi = int(m.group(1)), int(m.group(2))
        if (lo, hi) != (2, 3):
            add_flag(flags, 'Key fact mismatch', 'medium', f'DRS at {lo}-{hi} months', 'Guideline specifies the DRS should be performed at 2-3 months post injury for traumatic VS/UWS.', '5', context(norm, m.start(), m.end()))

    for m in re.finditer(r'\bmri\b[^.]{0,80}?(\d{1,2})\s*(?:-|to)\s*(\d{1,2})\s*week', norm):
        lo, hi = int(m.group(1)), int(m.group(2))
        if (lo, hi) != (6, 8):
            add_flag(flags, 'Key fact mismatch', 'medium', f'MRI at {lo}-{hi} weeks', 'Guideline specifies MRI should be performed 6-8 weeks post injury in traumatic VS/UWS.', '5', context(norm, m.start(), m.end()))

    for m in re.finditer(r'\bspect\b[^.]{0,80}?(\d{1,2})\s*(?:-|to)\s*(\d{1,2})\s*month', norm):
        lo, hi = int(m.group(1)), int(m.group(2))
        if (lo, hi) != (1, 2):
            add_flag(flags, 'Key fact mismatch', 'medium', f'SPECT at {lo}-{hi} months', 'Guideline specifies SPECT should be performed 1-2 months post injury in traumatic VS/UWS.', '5', context(norm, m.start(), m.end()))

    # --- Precise terminology/definition checks (Appendix A-style glossary) ---
    for m in re.finditer(r'\bcoma\b[^.]{0,100}?eyes?\s+(?:are\s+|remain(?:ing|ed)?\s+)?open', norm):
        add_flag(flags, 'Possible contradiction', 'high', 'coma ... eyes open', 'Coma is defined by no evidence of wakefulness, including eyes remaining continuously closed. Eyes opening is inconsistent with a coma diagnosis and suggests VS/UWS or a higher level of consciousness.', '7', context(norm, m.start(), m.end()))
    for m in re.finditer(r'eyes?\s+(?:are\s+|remain(?:ing|ed)?\s+)?open[^.]{0,100}?\bcoma\b', norm):
        add_flag(flags, 'Possible contradiction', 'high', 'eyes open ... coma', 'Coma is defined by no evidence of wakefulness, including eyes remaining continuously closed. Eyes opening is inconsistent with a coma diagnosis and suggests VS/UWS or a higher level of consciousness.', '7', context(norm, m.start(), m.end()))

    for m in re.finditer(r'(?:vegetative state|vs/uws|unresponsive wakefulness)[^.]{0,150}?(purposeful behavior|follow(?:s|ing)?\s+commands|command[- ]following)', norm):
        add_flag(flags, 'Possible contradiction', 'high', 'VS/UWS ... purposeful behavior/commands', 'VS/UWS is defined by NO evidence of purposeful behavior. Command following or purposeful behavior indicates at least MCS, not VS/UWS.', '7', context(norm, m.start(), m.end()))

    def sentence_window(t, start, end):
        lo = t.rfind('.', 0, start)
        lo = 0 if lo == -1 else lo + 1
        hi = t.find('.', end)
        hi = len(t) if hi == -1 else hi + 1
        return t[lo:hi]

    def sign_present(window, patterns):
        for pat in patterns:
            for pm in re.finditer(pat, window):
                preceding = window[max(0, pm.start() - 25):pm.start()]
                if re.search(r'\b(?:no|not|without|absence of|lack(?:ing)? of|negative for)\s*$', preceding):
                    continue
                return True
        return False

    mcs_plus_patterns = [r'command[- ]following', r'follow(?:s|ing)?\s+commands?', r'intelligible speech']
    mcs_minus_patterns = [r'automatic movements?', r'object manipulation', r'localizing', r'visual pursuit', r'visual fixation', r'affective behaviors?']

    for m in re.finditer(r'mcs\+', norm):
        win = sentence_window(norm, m.start(), m.end())
        has_plus_sign = sign_present(win, mcs_plus_patterns)
        has_minus_sign = sign_present(win, mcs_minus_patterns)
        if has_minus_sign and not has_plus_sign:
            add_flag(flags, 'Possible contradiction', 'medium', 'MCS+ near MCS- behaviors', 'MCS+ requires behavioral evidence of preserved receptive language (e.g., command following, intelligible speech). The nearby behaviors described (e.g., automatic movements, object manipulation, visual pursuit) define MCS-, not MCS+.', '7', context(norm, m.start(), m.end()))

    for m in re.finditer(r'mcs-(?!\w)', norm):
        win = sentence_window(norm, m.start(), m.end())
        has_minus_sign = sign_present(win, mcs_minus_patterns)
        has_plus_sign = sign_present(win, mcs_plus_patterns)
        if has_plus_sign and not has_minus_sign:
            add_flag(flags, 'Possible contradiction', 'medium', 'MCS- near MCS+ behaviors', 'MCS- is defined by nonlinguistic signs only (automatic movements, object manipulation, visual pursuit/fixation, affective behaviors). Command following or intelligible speech nearby indicates MCS+, not MCS-.', '7', context(norm, m.start(), m.end()))

    for m in re.finditer(r'persistent vegetative state[^.]{0,80}?(irreversible|permanent)', norm):
        add_flag(flags, 'Terminology', 'medium', 'persistent vegetative state ... permanent/irreversible', '"Persistent vegetative state" (PVS) denotes VS/UWS lasting more than 1 month and does not itself imply irreversibility. "Permanent vegetative state" is a distinct prognostic term applied at 3 months (nontraumatic) or 12 months (traumatic) indicating high probability of irreversibility.', '7', context(norm, m.start(), m.end()))

    for m in re.finditer(r'locked-in syndrome[^.]{0,150}?vegetative state', norm):
        add_flag(flags, 'Terminology', 'high', 'locked-in syndrome ... vegetative state', 'Locked-in syndrome (tetraplegia, anarthria, near-normal cognition) is a distinct condition that can be misdiagnosed as VS/UWS. It is not itself a disorder of consciousness and should not be equated with vegetative state.', None, context(norm, m.start(), m.end()))

    # Untrained clinicians / non-standardized assessments — contradicts Recs 1 & 4 (specialized,
    # trained multidisciplinary teams performing STANDARDIZED behavioral evaluations).
    for m in re.finditer(r'untrained clinicians?', norm):
        win = sentence_window(norm, m.start(), m.end())
        if re.search(r'non[\s-]?standardized|unstandardized', win):
            add_flag(flags, 'Possible contradiction', 'high', 'untrained clinicians ... non-standardized assessment',
                      'Guideline Recommendations 1 and 4 call for specialized, trained multidisciplinary teams to perform serial STANDARDIZED '
                      'behavioral evaluations — the opposite of untrained clinicians using non-standardized assessments.',
                      '4', context(norm, m.start(), m.end()))
    for m in re.finditer(r'non[\s-]?standardized\s+assessments?', norm):
        win = sentence_window(norm, m.start(), m.end())
        if not re.search(r'untrained clinicians?', win):  # avoid double-flagging the same sentence
            add_flag(flags, 'Possible contradiction', 'high', 'non-standardized assessment recommended',
                      'Guideline Recommendation 4 calls for performing serial STANDARDIZED behavioral evaluations to establish prognosis (Level B) — '
                      'not non-standardized assessments.',
                      '4', context(norm, m.start(), m.end()))

    # Premature withdrawal of life-sustaining treatment before the 28-day window — contradicts Rec 3
    # (must avoid statements suggesting a universally poor prognosis during the first 28 days, Level A).
    for m in re.finditer(r'withdrawal of life[\s-]?sustaining treatments?|withdraw(?:ing)? life[\s-]?sustaining treatments?', norm):
        win = sentence_window(norm, m.start(), m.end())
        time_match = re.search(r'(\d{1,3})\s*[- ]?\s*(hour|day)s?\b', win)
        if time_match:
            qty = int(time_match.group(1))
            unit = time_match.group(2)
            hours = qty if unit == 'hour' else qty * 24
            if hours < 672:  # 28 days
                add_flag(flags, 'Possible contradiction', 'high', f'withdrawal of life-sustaining treatment ... {qty} {unit}(s)',
                          f'Guideline Recommendation 3 states clinicians MUST avoid statements suggesting a universally poor prognosis during the '
                          f'first 28 days post injury (Level A). Recommending withdrawal of life-sustaining treatment based on early appearance, '
                          f'well before the 28-day window ({qty} {unit}(s) here), directly contradicts this.',
                          '3', context(norm, m.start(), m.end()))
    for m in re.finditer(r'vegetative state[^.]{0,150}?locked-in syndrome', norm):
        add_flag(flags, 'Terminology', 'high', 'vegetative state ... locked-in syndrome', 'Locked-in syndrome (tetraplegia, anarthria, near-normal cognition) is a distinct condition that can be misdiagnosed as VS/UWS. It is not itself a disorder of consciousness and should not be equated with vegetative state.', None, context(norm, m.start(), m.end()))

    for m in re.finditer(r'emerg(?:ed|ence)\s+from\s+mcs[^.]{0,150}?(?:used|using|use of)\s+(?:a|one|1)\s+(?:single\s+)?(?:familiar\s+)?object', norm):
        add_flag(flags, 'Key fact mismatch', 'medium', 'emergence from MCS ... one object', 'Emergence from MCS (EMCS) via functional object use requires demonstrated use of at least 2 different familiar objects, not a single object.', '7', context(norm, m.start(), m.end()))

    for m in re.finditer(r'(\d{1,2})\s*hours?\s+of\s+(?:multidisciplinary\s+)?therapy\s+(?:daily|per\s+day)', norm):
        hours = int(m.group(1))
        if hours != 3:
            add_flag(flags, 'Key fact mismatch', 'high', f'{hours} hours of therapy daily',
                      f'Standard IRF admission criteria require patients be able to tolerate approximately 3 hours of multidisciplinary therapy daily (the "3-hour rule"), not {hours} hours.',
                      '1', context(norm, m.start(), m.end()))

    for m in re.finditer(r'physicians?\s*\([^)]*intensivist[^)]*\)\.?\s*(?:the\s+)?overall leaders?', norm):
        add_flag(flags, 'Key fact mismatch', 'medium', 'Physicians ... overall leaders',
                  'The physician role (Intensivist, Neurologist, Physiatrist) on a multidisciplinary DoC team is typically described as team oversight, medical management, and disposition planning — not simply "the overall leaders."',
                  '1', context(norm, m.start(), m.end()))

    # --- Negated / inverted recommendation detector ---
    # Catches statements that flip the polarity of a guideline recommendation
    # (e.g., "referral ... is not critical" when the guideline recommends it).
    negation_phrase = r'not\s+(?:critical|necessary|essential|important|required|needed|recommended|beneficial|effective)|unnecessary|no\s+(?:need|benefit|role)|does\s+not\s+(?:improve|help|benefit)|isn.t\s+(?:critical|necessary|essential|important|required|needed)'
    negated_rec_topics = [
        {'keywords': [r'referral to a specialized rehabilitation', r'referral to (?:a |an )?multidisciplinary', r'multidisciplinary rehabilitation'], 'rec': '1',
         'note': 'Guideline Recommendation 1 states clinicians SHOULD refer medically stable patients with DoC to specialized multidisciplinary rehabilitation settings (Level B) to optimize diagnosis, prognostication, and management.'},
        {'keywords': [r'\bamantadine\b'], 'rec': '14',
         'note': 'Guideline Recommendation 14 supports amantadine (100-200 mg twice daily) to hasten functional recovery in appropriate traumatic VS/UWS or MCS patients (Level B).'},
        {'keywords': [r'goals of care'], 'rec': '9',
         'note': 'Guideline Recommendation 9 states clinicians MUST counsel families to establish goals of care once prognosis indicates likely severe long-term disability (Level A).'},
        {'keywords': [r'patient and family preferences', r'family preferences'], 'rec': '11',
         'note': 'Guideline Recommendation 11 requires incorporating patient and family preferences into care decisions (Level A).'},
        {'keywords': [r'pain assessment', r'pain management'], 'rec': '13',
         'note': 'The guideline supports routine pain assessment and management in patients with DoC (Level B).'},
        {'keywords': [r'serial standardized behavioral evaluations', r'serial behavioral evaluations'], 'rec': '4',
         'note': 'Guideline Recommendation 4 states clinicians SHOULD perform serial standardized behavioral evaluations to establish prognosis (Level B).'},
    ]
    for topic in negated_rec_topics:
        for kw in topic['keywords']:
            for m in re.finditer(kw + r'[^.]{0,150}?(?:' + negation_phrase + r')', norm):
                add_flag(flags, 'Possible inverted recommendation', 'high', m.group(0)[:120], f"This appears to invert the guideline's actual recommendation. {topic['note']}", topic['rec'], context(norm, m.start(), m.end()))
            for m in re.finditer(r'(?:' + negation_phrase + r')[^.]{0,150}?' + kw, norm):
                add_flag(flags, 'Possible inverted recommendation', 'high', m.group(0)[:120], f"This appears to invert the guideline's actual recommendation. {topic['note']}", topic['rec'], context(norm, m.start(), m.end()))

    topic_levels = {
        'amantadine': ('14', ['B']),
        'crs-r': ('6', ['B']),
        'coma recovery': ('6', ['B']),
        'sep': ('6', ['C']),
        'multidisciplinary rehabilitation': ('1', ['B']),
        'patient and family preferences': ('11', ['A']),
        'goals of care': ('9', ['A']),
        'pain': ('13', ['B']),
        'serial standardized': ('4', ['B']),
        'drs': ('5', ['B']),
        'spect': ('5', ['B']),
        'chronic phase': ('10', ['B'])
    }
    for m in re.finditer(r'Level\s+([ABCU])\b', text, re.IGNORECASE):
        cited = m.group(1).upper()
        win_lo = max(0, m.start() - 220)
        win = low[win_lo:m.end()+50]
        best = None
        for kw, (rid, levels) in topic_levels.items():
            pos = win.rfind(kw)
            if pos != -1:
                dist = abs((win_lo + pos) - m.start())
                if best is None or dist < best[0]:
                    best = (dist, kw, rid, levels)
        if best and cited not in best[3]:
            add_flag(flags, 'Evidence-level mismatch', 'high', f'Level {cited} near {best[1]}', f'Material cites Level {cited}, but guideline recommendation {best[2]} is Level {"/".join(best[3])}.', best[2], context(text, m.start(), m.end()))

    strong_hits = [t for t in STRONG_ANCHOR_TERMS if t in norm]
    weak_hits = [t for t in WEAK_ANCHOR_TERMS if t in norm]
    rec_sig_hit = any(re.search(p, norm) for patterns in REC_COVERAGE_SIGNATURES.values() for p in patterns)
    is_relevant = len(strong_hits) >= 1 or (len(strong_hits) + len(weak_hits)) >= 2 or rec_sig_hit

    coverage = []
    if is_relevant:
        for rec in kb['recommendations']:
            patterns = REC_COVERAGE_SIGNATURES.get(rec['id'], [])
            if any(re.search(p, norm) for p in patterns):
                coverage.append(rec['id'])
    else:
        add_flag(
            flags, 'Document relevance', 'high', '(no guideline-related terms detected)',
            f"This document does not appear to be about the {kb['meta']['title']} guideline — very few or no related clinical terms "
            "(e.g., vegetative state, minimally conscious state, disorders of consciousness, CRS-R, TBI) were found. The recommendation "
            "coverage map and fact-check results below are not meaningful for unrelated content and have been suppressed.",
            None, '',
        )

    quiz_spans = detect_quiz_spans(norm)
    if quiz_spans:
        surviving_flags = []
        for f in flags:
            pos = _flag_position(norm, f)
            in_quiz = pos != -1 and any(s <= pos < e for s, e, _ in quiz_spans)
            if in_quiz:
                continue  # quiz stems/distractors reference terms as their subject, not as assertions
            surviving_flags.append(f)
        flags = surviving_flags

        for start, end, answer_letter in quiz_spans:
            block_text = norm[start:end]
            for vf in verify_quiz_answer(block_text, answer_letter, kb):
                add_flag(flags, 'Possible quiz answer error', 'high', f'quiz answer "{answer_letter}"', vf['issue'], vf.get('rec'), context(norm, start, end))

    order = {'high':0, 'medium':1, 'low':2}
    seen = set()
    deduped_flags = []
    for f in flags:
        key = (f['kind'], f['severity'], f['matched'], f['issue'])
        if key in seen:
            continue
        seen.add(key)
        deduped_flags.append(f)
    flags = deduped_flags
    flags.sort(key=lambda x: order.get(x['severity'], 9))
    accessibility_findings = check_accessibility(filepath, text)
    return {
        'filename': os.path.basename(filepath),
        'text_length': len(text),
        'flags': flags,
        'coverage': sorted(set(coverage)),
        'clarity': analyze_clarity(text),
        'accessibility': accessibility_findings,
        'accessibility_score': compute_accessibility_score(accessibility_findings),
        'proofreading': analyze_proofreading(text, kb),
        'ocr_warning': ocr_warning,
    }


def make_report(result, kb, out_path):
    doc = Document()
    doc.add_heading('DoC Guideline Fact-Check Report', level=0)
    doc.add_paragraph(f"Material checked: {result['filename']}")
    doc.add_paragraph(f"Guideline: {kb['meta']['title']} ({kb['meta']['year']})")
    doc.add_paragraph(kb['meta']['citation'])
    doc.add_paragraph('This is a human-in-the-loop review aid. Flags identify candidates for review, not final clinical judgments.')

    counts = {s: 0 for s in ['high', 'medium', 'low']}
    for f in result['flags']:
        counts[f['severity']] = counts.get(f['severity'], 0) + 1
    doc.add_heading('Summary', level=1)
    doc.add_paragraph(f"High: {counts['high']} | Medium: {counts['medium']} | Low: {counts['low']}")

    doc.add_heading('Flags', level=1)
    if not result['flags']:
        doc.add_paragraph('No flags raised. Human review is still recommended.')
    colors = {'high': RGBColor(192,0,0), 'medium': RGBColor(199,106,0), 'low': RGBColor(127,106,0)}
    for f in result['flags']:
        p = doc.add_paragraph(style='List Bullet')
        r = p.add_run(f"[{f['severity'].upper()}] {f['kind']}")
        r.bold = True; r.font.color.rgb = colors.get(f['severity'], RGBColor(0,0,0))
        p.add_run(f" — matched: {f['matched']}")
        if f['rec']:
            p.add_run(f" (Rec {f['rec']})")
        doc.add_paragraph('Issue: ' + f['issue'])
        if f['context']:
            c = doc.add_paragraph('Context: ' + f['context'])
            c.runs[0].italic = True

    doc.add_heading('Recommendation Coverage Map', level=1)
    table = doc.add_table(rows=1, cols=4)
    table.style = 'Light Grid Accent 1'
    hdr = table.rows[0].cells
    hdr[0].text = 'Rec'; hdr[1].text = 'Topic'; hdr[2].text = 'Level'; hdr[3].text = 'Touched?'
    for rec in kb['recommendations']:
        row = table.add_row().cells
        row[0].text = rec['id']; row[1].text = rec['topic']; row[2].text = '/'.join(rec['level']); row[3].text = 'Yes' if rec['id'] in result['coverage'] else ''

    doc.add_heading('Appendix: Guideline Recommendation Wording', level=1)
    for rec in kb['recommendations']:
        p = doc.add_paragraph()
        p.add_run(f"Recommendation {rec['id']} (Level {'/'.join(rec['level'])}): ").bold = True
        p.add_run(rec['text'])
    doc.save(out_path)
